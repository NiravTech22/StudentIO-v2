"""
Document Parser
Extracts text from various file formats (PDF, DOCX, images, PPTX)
"""

import os
from typing import Dict, List, Any, Optional
import tempfile

# Document processing libraries
try:
    import PyPDF2
    import pdfplumber
except ImportError:
    PyPDF2 = None
    pdfplumber = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    import pytesseract
    from PIL import Image
except ImportError:
    pytesseract = None
    Image = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


class DocumentParser:
    """Handles parsing of various document formats"""
    
    @staticmethod
    def parse_pdf(file_path: str) -> Dict[str, Any]:
        """Extract text from PDF file"""
        if not PyPDF2 or not pdfplumber:
            raise ImportError("PDF libraries not installed. Install PyPDF2 and pdfplumber.")
        
        text_content = []
        metadata = {"pages": 0, "images": 0}
        
        try:
            # Try pdfplumber first (better for tables and complex layouts)
            with pdfplumber.open(file_path) as pdf:
                metadata["pages"] = len(pdf.pages)
                
                for page_num, page in enumerate(pdf.pages, 1):
                    page_text = page.extract_text()
                    
                    if page_text:
                        text_content.append({
                            "page": page_num,
                            "content": page_text,
                            "type": "text"
                        })
                   
                    # Extract tables
                    tables = page.extract_tables()
                    if tables:
                        for table_idx, table in enumerate(tables):
                            table_text = "\n".join([" | ".join([str(cell) for cell in row]) for row in table])
                            text_content.append({
                                "page": page_num,
                                "content": table_text,
                                "type": "table",
                                "table_index": table_idx
                            })
        
        except Exception as e:
            print(f"pdfplumber failed, falling back to PyPDF2: {e}")
            
            # Fallback to PyPDF2
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                metadata["pages"] = len(reader.pages)
                
                for page_num, page in enumerate(reader.pages, 1):
                    page_text = page.extract_text()
                    if page_text:
                        text_content.append({
                            "page": page_num,
                            "content": page_text,
                            "type": "text"
                        })
        
        return {
            "content": text_content,
            "metadata": metadata,
            "file_type": "pdf"
        }
    
    @staticmethod
    def parse_docx(file_path: str) -> Dict[str, Any]:
        """Extract text from DOCX file"""
        if not Document:
            raise ImportError("python-docx not installed")
        
        doc = Document(file_path)
        
        text_content = []
        
        # Extract paragraphs
        for para_idx, para in enumerate(doc.paragraphs):
            if para.text.strip():
                text_content.append({
                    "paragraph": para_idx + 1,
                    "content": para.text,
                    "type": "paragraph",
                    "style": para.style.name
                })
        
        # Extract tables
        for table_idx, table in enumerate(doc.tables):
            table_text = []
            for row in table.rows:
                row_text = [cell.text for cell in row.cells]
                table_text.append(" | ".join(row_text))
            
            text_content.append({
                "table": table_idx + 1,
                "content": "\n".join(table_text),
                "type": "table"
            })
        
        metadata = {
            "paragraphs": len(doc.paragraphs),
            "tables": len(doc.tables)
        }
        
        return {
            "content": text_content,
            "metadata": metadata,
            "file_type": "docx"
        }
    
    @staticmethod
    def parse_image(file_path: str, lang: str = 'eng') -> Dict[str, Any]:
        """Extract text from image using OCR"""
        if not pytesseract or not Image:
            raise ImportError("pytesseract and PIL not installed")
        
        try:
            # Open and process image
            img = Image.open(file_path)
            
            # Perform OCR
            text = pytesseract.image_to_string(img, lang=lang)
            
            # Get confidence data
            data = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT, lang=lang)
            
            # Calculate average confidence
            confidences = [float(conf) for conf in data['conf'] if conf != '-1']
            avg_confidence = sum(confidences) / len(confidences) if confidences else 0
            
            metadata = {
                "size": img.size,
                "mode": img.mode,
                "avg_confidence": avg_confidence,
                "language": lang
            }
            
            return {
                "content": [{"content": text, "type": "ocr"}],
                "metadata": metadata,
                "file_type": "image"
            }
        
        except Exception as e:
            raise Exception(f"OCR failed: {e}")
    
    @staticmethod
    def parse_pptx(file_path: str) -> Dict[str, Any]:
        """Extract text from PowerPoint presentation"""
        if not Presentation:
            raise ImportError("python-pptx not installed")
        
        prs = Presentation(file_path)
        
        text_content = []
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_text = []
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if slide_text:
                text_content.append({
                    "slide": slide_idx,
                    "content": "\n".join(slide_text),
                    "type": "slide"
                })
        
        metadata = {
            "slides": len(prs.slides),
            "slide_width": prs.slide_width,
            "slide_height": prs.slide_height
        }
        
        return {
            "content": text_content,
            "metadata": metadata,
            "file_type": "pptx"
        }
    
    @staticmethod    def parse_file(file_path: str, file_type: Optional[str] = None) -> Dict[str, Any]:
        """
        Auto-detect and parse any supported file type
        
        Args:
            file_path: Path to file
            file_type: Optional explicit file type
            
        Returns:
            Dict with 'content', 'metadata', and 'file_type'
        """
        # Detect file type if not provided
        if not file_type:
            ext = os.path.splitext(file_path)[1].lower()
            file_type = ext[1:] if ext else None
        
        # Route to appropriate parser
        if file_type == 'pdf':
            return DocumentParser.parse_pdf(file_path)
        elif file_type in ['docx', 'doc']:
            return DocumentParser.parse_docx(file_path)
        elif file_type in ['png', 'jpg', 'jpeg', 'tiff', 'bmp']:
            return DocumentParser.parse_image(file_path)
        elif file_type in ['pptx', 'ppt']:
            return DocumentParser.parse_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
    
    @staticmethod
    def extract_plain_text(parsed_doc: Dict[str, Any]) -> str:
        """Extract all text content as a single string"""
        content_items = parsed_doc.get("content", [])
        
        text_parts = []
        for item in content_items:
            if isinstance(item, dict):
                text_parts.append(item.get("content", ""))
            else:
                text_parts.append(str(item))
        
        return "\n\n".join(text_parts)
